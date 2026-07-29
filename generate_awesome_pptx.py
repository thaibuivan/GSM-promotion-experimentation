import collections
import collections.abc
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- COLORS ---
BG_COLOR = RGBColor(11, 19, 32)
CARD_BG = RGBColor(23, 37, 69)
YELLOW = RGBColor(255, 184, 0)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(52, 211, 153)
PINK = RGBColor(244, 114, 182)
WHITE = RGBColor(255, 255, 255)
FOOTER_BLUE = RGBColor(28, 78, 128)
FOOTER_DARK = RGBColor(7, 13, 23)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def create_base_layout(slide, title_text, subtitle_text, slide_num, total_slides):
    apply_background(slide)
    
    # Border
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0), Inches(0.15), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = YELLOW
    left_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.2), Inches(8.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(8.5), Inches(0.4))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = YELLOW
    
    # Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(2.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = YELLOW
    line.line.fill.background()
    
    # Footers
    footer1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(6.8), Inches(9.6), Inches(0.4))
    footer1.fill.solid()
    footer1.fill.fore_color.rgb = FOOTER_BLUE
    footer1.line.fill.background()
    
    f1_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(8), Inches(0.4))
    tf_f1 = f1_text.text_frame
    tf_f1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_f1 = tf_f1.paragraphs[0]
    p_f1.text = "GSM Promotion Experimentation • Báo cáo Chuyên môn • Data Science Team"
    p_f1.font.name = 'Segoe UI'
    p_f1.font.size = Pt(10)
    p_f1.font.color.rgb = WHITE
    
    footer2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(7.2), Inches(9.6), Inches(0.3))
    footer2.fill.solid()
    footer2.fill.fore_color.rgb = FOOTER_DARK
    footer2.line.fill.background()
    
    f2_text = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(5), Inches(0.3))
    tf_f2 = f2_text.text_frame
    tf_f2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_f2 = tf_f2.paragraphs[0]
    p_f2.text = "Project: Uplift Modeling & Causal Inference"
    p_f2.font.name = 'Segoe UI'
    p_f2.font.size = Pt(9)
    p_f2.font.color.rgb = CYAN
    
    num_text = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1.4), Inches(0.3))
    tf_num = num_text.text_frame
    tf_num.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_num = tf_num.paragraphs[0]
    p_num.text = f"{slide_num:02d} / {total_slides:02d}"
    p_num.font.name = 'Segoe UI'
    p_num.font.size = Pt(9)
    p_num.font.color.rgb = WHITE
    p_num.font.bold = True
    p_num.alignment = PP_ALIGN.RIGHT

def add_cards(slide, cards, image_path=None):
    num_cards = len(cards)
    if num_cards == 0: return
    
    # If there is an image, we split into 2 columns: Left for cards, Right for Image
    if image_path and os.path.exists(image_path):
        card_width = Inches(4.2)
        start_x = Inches(0.7)
        y_pos = Inches(1.6)
        card_height = Inches(4.9 / num_cards) # Stack them vertically
        spacing = Inches(0.15)
        
        for i, card in enumerate(cards):
            current_y = y_pos + i * (card_height + spacing)
            
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, current_y, card_width, card_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = CARD_BG
            rect.line.fill.background()
            
            accent_color = card.get('accent', CYAN)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x + Inches(0.15), current_y + Inches(0.4), card_width - Inches(0.3), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = accent_color
            line.line.fill.background()
            
            title_box = slide.shapes.add_textbox(start_x + Inches(0.05), current_y + Inches(0.05), card_width - Inches(0.1), Inches(0.4))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = card['title']
            p.font.name = 'Segoe UI'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = accent_color
            
            content_box = slide.shapes.add_textbox(start_x + Inches(0.05), current_y + Inches(0.45), card_width - Inches(0.1), card_height - Inches(0.5))
            tf_c = content_box.text_frame
            tf_c.word_wrap = True
            
            lines = card['content'].split('\n')
            for idx, text in enumerate(lines):
                if not text.strip(): continue
                p = tf_c.paragraphs[0] if idx == 0 else tf_c.add_paragraph()
                p.text = "• " + text
                p.font.name = 'Segoe UI'
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                p.space_after = Pt(4)
                
        # Add the image on the right (centered vertically)
        img_left = Inches(5.0)
        img_top = Inches(2.5)
        img_width = Inches(4.6)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        
    else:
        # Standard horizontal layout
        card_width = Inches(8.8 / num_cards)
        spacing = Inches(0.2)
        start_x = Inches(0.7)
        y_pos = Inches(1.8)
        card_height = Inches(4.5)
        
        for i, card in enumerate(cards):
            x_pos = start_x + i * (card_width + spacing)
            
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, card_width, card_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = CARD_BG
            rect.line.fill.background()
            
            accent_color = card.get('accent', CYAN)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.15), y_pos + Inches(0.6), card_width - Inches(0.3), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = accent_color
            line.line.fill.background()
            
            title_box = slide.shapes.add_textbox(x_pos + Inches(0.05), y_pos + Inches(0.05), card_width - Inches(0.1), Inches(0.55))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = card['title']
            p.font.name = 'Segoe UI'
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = accent_color
            
            content_box = slide.shapes.add_textbox(x_pos + Inches(0.05), y_pos + Inches(0.7), card_width - Inches(0.1), card_height - Inches(0.8))
            tf_c = content_box.text_frame
            tf_c.word_wrap = True
            
            lines = card['content'].split('\n')
            for idx, text in enumerate(lines):
                if not text.strip(): continue
                p = tf_c.paragraphs[0] if idx == 0 else tf_c.add_paragraph()
                p.text = "• " + text
                p.font.name = 'Segoe UI'
                p.font.size = Pt(15)
                p.font.color.rgb = WHITE
                p.space_after = Pt(6)

def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9.0), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Segoe UI'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(4.2), Inches(2.0), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = YELLOW
    line.line.fill.background()
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9.0), Inches(1.0))
    tf_sub = sub_box.text_frame
    tf_sub.vertical_anchor = MSO_ANCHOR.TOP
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = CYAN

def create_week7_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    TOTAL_SLIDES = 6
    create_title_slide(prs, "BÁO CÁO TUẦN 7: UPLIFT MODELING", "Tối ưu hóa Lợi nhuận bằng Causal AI")
    
    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, "1. Tại sao cần Uplift Modeling?", "Giải bài toán Ăn thịt Doanh thu", 1, TOTAL_SLIDES)
    add_cards(s1, [
        {'title': 'Hiệu ứng Cannibalization', 'content': 'Khách hàng "Sure Things" đằng nào cũng mua.\nPhát mã cho nhóm này là ném tiền qua cửa sổ.', 'accent': YELLOW},
        {'title': 'Giải pháp Uplift', 'content': 'Dự đoán CATE cho TỪNG người.\nTìm ra nhóm "Persuadables" - Chỉ mua khi có mã.', 'accent': CYAN}
    ])
    
    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, "2. Thiết kế & Phân phối CATE", "Sự đa dạng trong tác động của Voucher", 2, TOTAL_SLIDES)
    add_cards(s2, [
        {'title': 'T-Learner (XGBoost)', 'content': 'Áp dụng Two-Model Approach để tìm CATE.\nÁp dụng Early Stopping chống Overfitting.', 'accent': CYAN},
        {'title': 'Biểu đồ Phân phối', 'content': 'Biểu đồ bên phải cho thấy tác động của Voucher có tính phân hóa rất rõ ràng ở từng Users.', 'accent': GREEN}
    ], image_path='docs/images/week7_uplift_modeling_0.png')

    # Slide 3
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, "3. Đánh giá Mô hình (Qini Curve)", "Mô hình vượt trội so với Random", 3, TOTAL_SLIDES)
    add_cards(s3, [
        {'title': 'Đường cong Qini', 'content': 'Chỉ số đo lường mức độ gia tăng chuyến đi lũy kế.\nĐường màu xanh vồng lên vượt trội so với đường chéo (Random).', 'accent': CYAN},
        {'title': 'Kết luận Đánh giá', 'content': 'Mô hình T-Learner hoàn toàn có khả năng xếp hạng (Ranking) độ nhạy cảm của khách.', 'accent': YELLOW}
    ], image_path='docs/images/week7_uplift_modeling_2.png')

    # Slide 4
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, "4. Chuyển hóa Mô hình thành Tiền", "Action Distribution", 4, TOTAL_SLIDES)
    add_cards(s4, [
        {'title': 'Quyết định phát mã', 'content': 'Mô hình cắt lọc và tính toán ra Lợi nhuận kỳ vọng.\nChỉ phát Voucher cho nhóm có Expected Profit > 0.', 'accent': CYAN},
        {'title': 'Tỷ lệ Action', 'content': 'Biểu đồ cho thấy tỷ lệ cắt giảm đáng kể số lượng Voucher không hiệu quả.', 'accent': PINK}
    ], image_path='docs/images/week7_uplift_modeling_1.png')

    # Slide 5
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s5, "5. Ngưỡng Hòa Vốn (Break-even)", "Điểm sinh lời", 5, TOTAL_SLIDES)
    add_cards(s5, [
        {'title': 'Công thức Hòa Vốn', 'content': 'Chi phí Voucher: 15k, Lãi gộp: 20k.\nNgưỡng CATE sinh lời = 15k / 20k = 0.75 chuyến.', 'accent': YELLOW},
        {'title': 'Sự thật tàn nhẫn', 'content': 'Toàn bộ tập Urban Cash không có ai đạt mức CATE > 0.75.', 'accent': PINK}
    ])

    # Slide 6
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s6, "6. Đề xuất Hành động (Decision)", "Hướng đi tiếp theo", 6, TOTAL_SLIDES)
    add_cards(s6, [
        {'title': '1. DỪNG CHIẾN DỊCH URBAN', 'content': 'Dừng phát mã cho nhóm Urban Cash để bảo toàn vốn.', 'accent': YELLOW},
        {'title': '2. GATEKEEPER', 'content': 'Dùng T-Learner làm màng lọc CATE > 0.75 cho mọi campaign tương lai.', 'accent': CYAN},
        {'title': '3. DỊCH CHUYỂN NGÂN SÁCH', 'content': 'Tái phân bổ ngân sách sang nhóm Suburban tiềm năng hơn.', 'accent': GREEN}
    ])
    
    prs.save('docs/Week7_Uplift_Modeling_Summary.pptx')

def create_week8_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    TOTAL_SLIDES = 6
    create_title_slide(prs, "BÁO CÁO TUẦN 8: STRESS TESTING", "Kiểm định Tính vững (Robustness Check) của Hệ thống")
    
    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, "1. Tổng quan Stress Test", "Mục tiêu và Phương pháp tiếp cận", 1, TOTAL_SLIDES)
    add_cards(s1, [
        {'title': 'Mục tiêu Kiểm định', 'content': 'Kiểm tra độ tin cậy của mô hình Causal Inference trước khi triển khai thực tế.\nĐảm bảo hệ thống không bị sai lệch bởi các yếu tố ngoại cảnh hoặc quy mô dữ liệu.', 'accent': CYAN},
        {'title': 'Phương pháp', 'content': 'Thực hiện 4 bài kiểm tra độc lập (Scalability, A/A Test, Imbalanced Ratio, Noise Injection).\nĐo lường sự biến động của ATE và P-value trong từng kịch bản.', 'accent': YELLOW}
    ])
    
    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, "2. Bài test 1: Độ Hội Tụ (Scalability)", "Đánh giá sự ổn định khi tăng quy mô mẫu", 2, TOTAL_SLIDES)
    add_cards(s2, [
        {'title': 'Thiết lập', 'content': 'Tăng dần kích thước mẫu (Sample Size) từ 10,000 lên hàng trăm ngàn quan sát.', 'accent': CYAN},
        {'title': 'Kết quả', 'content': 'Estimated ATE bám sát True ATE.\nKhoảng tin cậy (95% CI) thu hẹp dần theo đúng định lý giới hạn trung tâm.', 'accent': GREEN}
    ], image_path='docs/images/week8_stress_test_3.png')

    # Slide 3
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, "3. Bài test 2: A/A Test (Giả dược)", "Kiểm định tỷ lệ Dương tính giả (False Positives)", 3, TOTAL_SLIDES)
    add_cards(s3, [
        {'title': 'Thiết lập Kịch bản', 'content': 'Giả lập tình huống Voucher không có bất kỳ tác động nào (True Effect = 0).\nChia ngẫu nhiên tập khách hàng làm 2 nhóm và so sánh.', 'accent': YELLOW},
        {'title': 'Kết quả Đo lường', 'content': 'P-value duy trì ở mức > 0.05.\nHệ thống kết luận chính xác không có sự khác biệt có ý nghĩa thống kê.', 'accent': GREEN},
        {'title': 'Đánh giá', 'content': 'Mô hình không bị "ảo giác" số liệu.\nLoại trừ rủi ro phân bổ ngân sách sai lầm do lỗi nhiễu ngẫu nhiên.', 'accent': CYAN}
    ])

    # Slide 4
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, "4. Bài test 3: Tỷ lệ Lệch Mẫu", "Đánh giá sức chịu đựng với kích thước mẫu nhỏ", 4, TOTAL_SLIDES)
    add_cards(s4, [
        {'title': 'Thiết lập Kịch bản', 'content': 'Giả định ngân sách giới hạn, buộc phải ép tỷ lệ chia mẫu Control/Treatment xuống mức 90/10.', 'accent': PINK},
        {'title': 'Kết quả Đo lường', 'content': 'Phương sai (Variance) tăng do mẫu Treatment mỏng đi.\nTuy nhiên, giá trị ATE ước lượng vẫn xoay quanh ATE thực tế.', 'accent': YELLOW},
        {'title': 'Đánh giá', 'content': 'Mô hình vẫn hoạt động ổn định ngay cả khi điều kiện thử nghiệm không lý tưởng (không thể chia 50/50).', 'accent': GREEN}
    ])

    # Slide 5
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s5, "5. Bài test 4: Kháng Nhiễu (Noise)", "Kiểm tra độ vững trước tác động ngoại cảnh", 5, TOTAL_SLIDES)
    add_cards(s5, [
        {'title': 'Thiết lập Kịch bản', 'content': 'Tiêm nhiễu ngẫu nhiên (Gaussian Noise) vào biến mục tiêu (Số chuyến đi).\nMô phỏng tác động của các sự kiện bất thường (ví dụ: thời tiết xấu, lễ hội).', 'accent': YELLOW},
        {'title': 'Kết quả Đo lường', 'content': 'Phân phối số chuyến đi bị biến dạng mạnh.\nTuy nhiên, chênh lệch ATE giữa 2 nhóm T và C gần như không thay đổi.', 'accent': CYAN},
        {'title': 'Cơ sở Toán học', 'content': 'Quá trình phân bổ ngẫu nhiên (Randomization) phân phối đều nhiễu lên cả 2 nhóm.\nKhi tính toán ATE (T trừ C), lượng nhiễu này tự triệt tiêu.', 'accent': GREEN}
    ])

    # Slide 6
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s6, "6. Tổng kết & Đề xuất Triển khai", "Đánh giá mức độ sẵn sàng của hệ thống", 6, TOTAL_SLIDES)
    add_cards(s6, [
        {'title': 'Đánh giá Tính vững (Robustness)', 'content': 'Hệ thống đo lường Causal Inference duy trì tính chính xác qua cả 4 kịch bản Stress Test tiêu chuẩn.', 'accent': GREEN},
        {'title': 'Khuyến nghị Triển khai', 'content': 'Đạt tiêu chuẩn kỹ thuật để tích hợp vào Pipeline dữ liệu thực tế.\nSẵn sàng triển khai A/B Testing trên quy mô lớn.', 'accent': CYAN}
    ])
    
    prs.save('docs/Week8_Stress_Test_Summary.pptx')

if __name__ == '__main__':
    create_week7_ppt()
    create_week8_ppt()
