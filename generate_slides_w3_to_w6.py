from pptx import Presentation
import sys
import os

sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from generate_awesome_pptx import (
    create_title_slide, create_base_layout, add_cards,
    YELLOW, CYAN, GREEN, PINK, WHITE
)

def create_week3_ppt():
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    TOTAL = 7
    
    create_title_slide(prs, "BÁO CÁO TUẦN 3: EXPERIMENT DESIGN", "Thiết kế Thử nghiệm Kích hoạt Lại Khách hàng Ngủ đông")
    
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, "1. Bối cảnh & Vấn đề (Rationale)", "Nguy cơ Rời bỏ của Regular Users", 1, TOTAL)
    add_cards(s1, [
        {'title': 'Tệp Khách hàng Giá trị', 'content': 'Regular Users (4-6 chuyến/tháng) là tệp có giá trị nhất nhưng dễ rơi vào trạng thái "Ngủ đông" sau 5-10 ngày không mở app.', 'accent': YELLOW},
        {'title': 'Vấn đề của Chiến dịch Cũ', 'content': 'Phát mã đại trà (Time-based) gây ra hiệu ứng Cannibalization (Tặng tiền cho người vốn đã định đặt xe).', 'accent': PINK},
        {'title': 'Giải pháp v2', 'content': 'Chuyển sang nhắm mục tiêu theo phân khúc (User-segment-based) để tối ưu hóa ngân sách và đo lường Causal Effect.', 'accent': CYAN}
    ])
    
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, "2. Giả thuyết & Estimand", "Định lượng Tác động", 2, TOTAL)
    add_cards(s2, [
        {'title': 'Giả thuyết', 'content': 'Gửi Push Notification kèm Voucher 25% sẽ phá vỡ "quán tính" của khách hàng, kích thích họ mở app và đặt xe trở lại.', 'accent': CYAN},
        {'title': 'ATE (Average Treatment Effect)', 'content': 'Sự chênh lệch số chuyến đi trung bình/user trong 14 ngày giữa nhóm nhận Voucher và nhóm Không nhận.', 'accent': GREEN}
    ])
    
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, "3. Tập Mục tiêu (Target Population)", "Lựa chọn Dựa trên K-Means", 3, TOTAL)
    add_cards(s3, [
        {'title': 'Tiêu chí Chọn vào (Inclusion)', 'content': 'Chỉ chọn nhóm Nhạy cảm giá (Persuadables): Urban Credit Card và Suburban Occasionals.\nĐã ngủ đông từ 5-14 ngày.', 'accent': CYAN},
        {'title': 'Tiêu chí Loại trừ (Exclusion)', 'content': 'Khách đi sân bay (Kháng sale) và Khách đi làm giờ cao điểm (Bắt buộc phải đi).\nKhách đã nhận Voucher trong 14 ngày qua.', 'accent': PINK}
    ], image_path='docs/images/week3_segmentation_0.png')
    
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, "4. Thiết kế Can thiệp", "Treatment vs Control", 4, TOTAL)
    add_cards(s4, [
        {'title': 'Nhóm Can thiệp (Treatment)', 'content': 'Hệ thống gửi Push Notification cá nhân hóa.\nHiển thị In-app Banner nhắc nhở.\nTự động áp mã giảm 25% khi thanh toán.', 'accent': GREEN},
        {'title': 'Nhóm Đối chứng (Control)', 'content': 'Không nhận thông báo, không có Banner.\nKhông chạy bất kỳ chiến dịch nào song song để tránh lây nhiễm.', 'accent': YELLOW}
    ])
    
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s5, "5. Hệ thống Chỉ số Đánh giá", "Metrics Framework", 5, TOTAL)
    add_cards(s5, [
        {'title': 'Chỉ số Cốt lõi (OEC)', 'content': 'Incremental Rides per User (Số chuyến đi tăng thêm trên mỗi User trong vòng 14 ngày).', 'accent': CYAN},
        {'title': 'Chỉ số Bảo vệ (Guardrails)', 'content': 'Profit Margin per Ride (Lợi nhuận không âm).\nUser Complaint Rate < 2% (Đảm bảo Push không bị đánh dấu Spam).', 'accent': YELLOW}
    ])
    
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s6, "6. Cỡ mẫu & Power Analysis", "Đảm bảo Ý nghĩa Thống kê", 6, TOTAL)
    add_cards(s6, [
        {'title': 'Thông số Thiết lập', 'content': 'Statistical Power: 80%.\nMức Ý nghĩa (Alpha): 5%.\nMinimum Detectable Effect (MDE): +0.05 chuyến/user.', 'accent': CYAN},
        {'title': 'Kết quả', 'content': 'Yêu cầu tối thiểu 2,500 users mỗi nhóm.\nTập dữ liệu hiện có 6,600 users, hoàn toàn vượt qua bài kiểm tra.', 'accent': GREEN}
    ])
    
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s7, "7. Biến Nhiễu & Luật Quyết định", "Bảo vệ Tính toàn vẹn của Thử nghiệm", 7, TOTAL)
    add_cards(s7, [
        {'title': 'Kiểm soát Confounders', 'content': 'Áp dụng Phân bổ ngẫu nhiên theo tầng (Stratified Randomization) để cân bằng các biến nhiễu như: Loại khách, Khu vực sống.', 'accent': YELLOW},
        {'title': 'Luật Ra Quyết định', 'content': 'Thí nghiệm thành công khi: P-value < 0.05 VÀ Incremental Rides > MDE VÀ Lợi nhuận Margin Dương.', 'accent': CYAN}
    ], image_path='docs/images/week3_segmentation_1.png')
    
    prs.save('docs/Week3_Experiment_Design_Summary.pptx')

def create_week4_ppt():
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    TOTAL = 5
    
    create_title_slide(prs, "BÁO CÁO TUẦN 4: A/B TESTING", "Nền tảng Thống kê & Phân tích Dữ liệu")
    
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, "1. Nền tảng Thống kê (Frequentist Testing)", "Kiểm định Giả thuyết", 1, TOTAL)
    add_cards(s1, [
        {'title': 'Giả thuyết Không (H0)', 'content': 'Giả định rằng Voucher không tạo ra sự khác biệt nào giữa 2 nhóm (Tác động = 0).', 'accent': YELLOW},
        {'title': 'P-value', 'content': 'Đại diện cho xác suất thu được chênh lệch do ngẫu nhiên.\nNgưỡng Alpha = 0.05 là tiêu chuẩn để bác bỏ H0.', 'accent': CYAN}
    ])
    
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, "2. Hai Tiêu chí Đánh giá Độc lập", "Sanity vs OEC", 2, TOTAL)
    add_cards(s2, [
        {'title': 'Sanity Checks', 'content': 'Kiểm tra độ cân bằng hệ thống (Sample Ratio, Covariates).\nKỳ vọng P-value > 0.05 (Không bác bỏ H0) để chứng minh hệ thống không bị lệch.', 'accent': PINK},
        {'title': 'OEC Evaluation', 'content': 'Đo lường tác động lên chỉ số cốt lõi (Incremental Rides).\nKỳ vọng P-value < 0.05 (Bác bỏ H0) để chứng minh Voucher có hiệu quả.', 'accent': GREEN}
    ])
    
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, "3. Đánh giá Tác động Ban đầu", "Chênh lệch giữa Control và Treatment", 3, TOTAL)
    add_cards(s3, [
        {'title': 'Quan sát trực quan', 'content': 'Biểu đồ cho thấy nhóm được nhận khuyến mãi (Treatment) có sự dịch chuyển tích cực về tần suất đặt xe.', 'accent': GREEN},
        {'title': 'Khẳng định', 'content': 'Bằng chứng ban đầu xác nhận Voucher thực sự tạo ra Tác động Nhân quả (Causal Effect).', 'accent': CYAN}
    ], image_path='docs/images/week4_ab_testing_0.png')
    
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, "4. Phân tích Hồi quy Tuyến tính (OLS)", "Đo lường Net Effect", 4, TOTAL)
    add_cards(s4, [
        {'title': 'Kiểm soát Covariates', 'content': 'Hồi quy OLS cho phép cô lập tác động của Voucher khỏi các biến nhiễu ngoại cảnh (Thời tiết, Vị trí).', 'accent': YELLOW},
        {'title': 'Tác động Ròng', 'content': 'Chỉ số Uplift đo lường được là chính xác tuyệt đối, loại bỏ sự may mắn ngẫu nhiên.', 'accent': CYAN}
    ], image_path='docs/images/week4_ab_testing_1.png')
    
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s5, "5. Khám phá Tác động Dị thể (HTE)", "Causal Trees", 5, TOTAL)
    add_cards(s5, [
        {'title': 'Heterogeneous Effect', 'content': 'Tác động của Voucher không phân bổ đều. Sẽ có nhóm cực kỳ nhạy cảm và nhóm hoàn toàn thờ ơ.', 'accent': PINK},
        {'title': 'Cây Quyết định Nhân quả', 'content': 'Thuật toán tự động tìm ra các nhánh (Segments) có CATE (Conditional ATE) cao nhất để nhắm mục tiêu.', 'accent': GREEN}
    ], image_path='docs/images/week4_ab_testing_2.png')
    
    prs.save('docs/Week4_AB_Testing_Summary.pptx')

def create_week5_ppt():
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    TOTAL = 4
    
    create_title_slide(prs, "BÁO CÁO TUẦN 5: A/A TESTING", "Kiểm định Độ tin cậy & Trustworthiness Checklist")
    
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, "1. Tầm quan trọng của A/A Test", "Sanity Check Cốt lõi", 1, TOTAL)
    add_cards(s1, [
        {'title': 'Môi trường Giả dược', 'content': 'Chạy thử nghiệm trên 2 nhóm hoàn toàn giống nhau (Không nhóm nào nhận Voucher).', 'accent': YELLOW},
        {'title': 'Nhiệm vụ', 'content': 'Kiểm chứng Engine phân bổ ngẫu nhiên (Randomization) hoạt động hoàn hảo, không tạo ra sự chênh lệch ảo.', 'accent': CYAN}
    ])
    
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, "2. Kiểm tra Covariate Balance", "Sự cân bằng của các Đặc tính", 2, TOTAL)
    add_cards(s2, [
        {'title': 'Phân phối Đồng đều', 'content': 'Các đặc tính (Tuổi, Giới tính, Khu vực, Lịch sử) giữa 2 nhóm A/A gần như trùng khớp hoàn toàn.', 'accent': GREEN},
        {'title': 'Đánh giá SRM', 'content': 'Hệ thống vượt qua kiểm định Sample Ratio Mismatch, tỷ lệ chia luôn duy trì ở mức 50.0/50.0.', 'accent': CYAN}
    ], image_path='docs/images/week5_aa_testing_0.png')
    
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, "3. Kháng lỗi Dương tính giả (False Positives)", "Mô phỏng Monte Carlo", 3, TOTAL)
    add_cards(s3, [
        {'title': '5.000 Vòng lặp', 'content': 'Thực hiện chạy A/A Test ngẫu nhiên 5.000 lần để kiểm tra độ ổn định của hệ thống.', 'accent': PINK},
        {'title': 'Tỷ lệ False Positive Rate', 'content': 'Tỷ lệ hệ thống kết luận nhầm "Có khác biệt" duy trì ở mức ~4.74% (Sát với ngưỡng lý tưởng 5%).', 'accent': GREEN}
    ])
    
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, "4. Tiêu chuẩn Trust Checklist", "Kết luận Độ tin cậy", 4, TOTAL)
    add_cards(s4, [
        {'title': 'P-value Uniformity', 'content': 'Giá trị P-value trải dài đều đặn (Uniform Distribution), chứng tỏ hệ thống đo lường không bị lệch (Bias).', 'accent': YELLOW},
        {'title': 'Green Light', 'content': 'Nền tảng đạt chuẩn quốc tế, hoàn toàn đủ tin cậy để triển khai các chiến dịch A/B Test tốn hàng tỷ đồng.', 'accent': GREEN}
    ])
    
    prs.save('docs/Week5_AA_Testing_Summary.pptx')

def create_week6_ppt():
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    TOTAL = 6
    
    create_title_slide(prs, "BÁO CÁO TUẦN 6: FINAL RCT RESULTS", "Sự bùng nổ của Hiệu ứng Dị thể (HTE)")
    
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, "1. Tóm tắt Thực thi (Executive Summary)", "Bức tranh Tổng thể", 1, TOTAL)
    add_cards(s1, [
        {'title': 'Sự thật Bất ngờ', 'content': 'Chiến dịch thử nghiệm đã vạch trần một sự thật quan trọng về Hiệu ứng Dị thể (Heterogeneous Treatment Effect).', 'accent': YELLOW},
        {'title': 'Thay đổi Chiến lược', 'content': 'Dừng ngay lập tức việc phát Voucher đại trà. Chuyển hướng toàn bộ ngân sách sang nhóm khách hàng thực sự sinh lời.', 'accent': GREEN}
    ])
    
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, "2. Định vị Phân khúc (Targeting)", "4 Cụm Hành vi K-Means", 2, TOTAL)
    add_cards(s2, [
        {'title': 'Nhóm Đi làm (Commuters)', 'content': 'Urban Commuters và Suburban Commuters: Có tính chất thiết yếu, cầu không co giãn với giá.', 'accent': CYAN},
        {'title': 'Nhóm Tự do (Occasionals)', 'content': 'Urban Credit Card và Suburban Occasionals: Khách đi chơi, độ nhạy cảm với khuyến mãi cực kỳ cao.', 'accent': PINK}
    ])
    
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, "3. Thiết kế Thí nghiệm (Experiment Design)", "Randomized Controlled Trial", 3, TOTAL)
    add_cards(s3, [
        {'title': 'Can thiệp (Treatment)', 'content': 'Tặng Voucher 25% cho một nửa số khách hàng được chọn ngẫu nhiên.', 'accent': CYAN},
        {'title': 'Đo lường OEC & Guardrail', 'content': 'OEC: Số chuyến đi gia tăng (Incremental Rides).\nGuardrail: Tỷ suất Hoàn vốn (ROI) chiến dịch > 0.', 'accent': YELLOW}
    ])
    
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, "4. Phân tích HTE: Cái bẫy Lợi nhuận", "Nhóm Suburban Commuters", 4, TOTAL)
    add_cards(s4, [
        {'title': 'Kết quả Đo lường', 'content': 'Số chuyến đi có tăng nhẹ (+0.84 chuyến, p < 0.05).', 'accent': YELLOW},
        {'title': 'Cannibalization (Ăn thịt doanh thu)', 'content': 'Việc trợ giá 25% cho những người đằng nào cũng phải đi làm dẫn đến doanh thu tăng thêm không bù nổi chi phí. ROI âm hoặc hòa vốn.', 'accent': PINK}
    ])
    
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s5, "5. Phân tích HTE: Mỏ vàng Thực sự", "Nhóm Urban Credit Card", 5, TOTAL)
    add_cards(s5, [
        {'title': 'Bùng nổ Nhu cầu', 'content': 'Khách nội thành đi chơi cực kỳ khoái khuyến mãi. Voucher đã kích hoạt số chuyến đi tăng vọt (+0.91 đến 2.5 chuyến).', 'accent': CYAN},
        {'title': 'Super ROI', 'content': 'Lợi nhuận ròng tăng đột biến, tỷ suất ROI vượt mốc 100%. Đây là tệp khách hàng "gánh" toàn bộ chiến dịch.', 'accent': GREEN}
    ])
    
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s6, "6. Lộ trình Triển khai Kế tiếp", "Hành động (Next Steps)", 6, TOTAL)
    add_cards(s6, [
        {'title': '1. Stop Mass Marketing', 'content': 'Chặn toàn bộ Voucher đối với nhóm Commuters để bảo toàn lợi nhuận biên (Profit Margin).', 'accent': YELLOW},
        {'title': '2. Roll-out', 'content': 'Chạy chiến dịch tự động tặng Voucher 25% cho tệp Urban Credit Card vào khung giờ thấp điểm.', 'accent': GREEN},
        {'title': '3. Causal ML (Uplift)', 'content': 'Bước tiếp theo: Dùng AI dự đoán chính xác từng cá nhân (Persuadables) để tối ưu hóa triệt để chi phí.', 'accent': CYAN}
    ])
    
    prs.save('docs/Week6_Final_Experiment_Summary.pptx')

if __name__ == '__main__':
    create_week3_ppt()
    create_week4_ppt()
    create_week5_ppt()
    create_week6_ppt()
