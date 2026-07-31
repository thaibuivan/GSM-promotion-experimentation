import sys
import os

sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from generate_awesome_pptx import (
    create_title_slide, create_base_layout, add_cards,
    YELLOW, CYAN, GREEN, PINK, WHITE
)

def create_merged_ppt():
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    TOTAL = 11  # 6 from W3 + 5 from W4
    current_slide = 1
    
    create_title_slide(prs, "BÁO CÁO TIẾN ĐỘ TUẦN 2", "K-Means Clustering & A/B Testing")
    
    # --- WEEK 3 CONTENT (SEGMENTATION) ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, f"{current_slide}. Tập Mục tiêu (Target Population)", "Lựa chọn Dựa trên K-Means", current_slide, TOTAL)
    add_cards(s3, [
        {'title': 'Tiêu chí Chọn vào (Inclusion)', 'content': 'Chỉ chọn nhóm Nhạy cảm giá (Persuadables): Urban Credit Card và Suburban Occasionals.\nĐã ngủ đông từ 5-14 ngày.', 'accent': CYAN},
        {'title': 'Tiêu chí Loại trừ (Exclusion)', 'content': 'Khách đi sân bay (Kháng sale) và Khách đi làm giờ cao điểm (Bắt buộc phải đi).\nKhách đã nhận Voucher trong 14 ngày qua.', 'accent': PINK}
    ], image_path='docs/images/week3_segmentation_0.png')
    current_slide += 1
    
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, f"{current_slide}. Thiết kế Can thiệp", "Treatment vs Control", current_slide, TOTAL)
    add_cards(s4, [
        {'title': 'Nhóm Can thiệp (Treatment)', 'content': 'Hệ thống gửi Push Notification cá nhân hóa.\nHiển thị In-app Banner nhắc nhở.\nTự động áp mã giảm 25% khi thanh toán.', 'accent': GREEN},
        {'title': 'Nhóm Đối chứng (Control)', 'content': 'Không nhận thông báo, không có Banner.\nKhông chạy bất kỳ chiến dịch nào song song để tránh lây nhiễm.', 'accent': YELLOW}
    ])
    current_slide += 1
    
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s5, f"{current_slide}. Hệ thống Chỉ số Đánh giá", "Metrics Framework", current_slide, TOTAL)
    add_cards(s5, [
        {'title': 'Chỉ số Cốt lõi (OEC)', 'content': 'Incremental Rides per User (Số chuyến đi tăng thêm trên mỗi User trong vòng 14 ngày).', 'accent': CYAN},
        {'title': 'Chỉ số Bảo vệ (Guardrails)', 'content': 'Profit Margin per Ride (Lợi nhuận không âm). Đảm bảo chiến dịch khuyến mãi không bị lỗ trên từng cuốc xe.', 'accent': YELLOW}
    ])
    current_slide += 1
    
    # --- WEEK 4 CONTENT (A/B TESTING) ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, f"{current_slide}. Nền tảng Thống kê", "Kiểm định Giả thuyết (Frequentist Testing)", current_slide, TOTAL)
    add_cards(s1, [
        {'title': 'Giả thuyết Không (H0)', 'content': 'Giả định rằng Voucher không tạo ra sự khác biệt nào giữa 2 nhóm (Tác động = 0).', 'accent': YELLOW},
        {'title': 'P-value', 'content': 'Đại diện cho xác suất thu được chênh lệch do ngẫu nhiên.\nNgưỡng Alpha = 0.05 là tiêu chuẩn để bác bỏ H0.', 'accent': CYAN}
    ])
    current_slide += 1
    
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, f"{current_slide}. Hai Tiêu chí Đánh giá Độc lập", "Sanity vs OEC", current_slide, TOTAL)
    add_cards(s2, [
        {'title': 'Sanity Checks', 'content': 'Kiểm tra độ cân bằng hệ thống (Sample Ratio, Covariates).\nKỳ vọng P-value > 0.05 (Không bác bỏ H0) để chứng minh hệ thống không bị lệch.', 'accent': PINK},
        {'title': 'OEC Evaluation', 'content': 'Đo lường tác động lên chỉ số cốt lõi (Incremental Rides).\nKỳ vọng P-value < 0.05 (Bác bỏ H0) để chứng minh Voucher có hiệu quả.', 'accent': GREEN}
    ])
    current_slide += 1
    
    s3_w4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3_w4, f"{current_slide}. Đánh giá Tác động Ban đầu", "Chênh lệch giữa Control và Treatment", current_slide, TOTAL)
    add_cards(s3_w4, [
        {'title': 'Quan sát trực quan', 'content': 'Biểu đồ cho thấy nhóm được nhận khuyến mãi (Treatment) có sự dịch chuyển tích cực về tần suất đặt xe.', 'accent': GREEN},
        {'title': 'Khẳng định', 'content': 'Bằng chứng ban đầu xác nhận Voucher thực sự tạo ra Tác động Nhân quả (Causal Effect).', 'accent': CYAN}
    ], image_path='docs/images/week4_ab_testing_0.png')
    current_slide += 1

    prs.save('docs/Progress_Update_Week2_Merged.pptx')
    print("Done generating docs/Progress_Update_Week2_Merged.pptx")

if __name__ == '__main__':
    create_merged_ppt()
