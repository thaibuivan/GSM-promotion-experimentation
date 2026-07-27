from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ======================================
# COLOUR PALETTE
# ======================================
CLR_BG       = RGBColor(0x0F, 0x1B, 0x2D)   # Deep navy
CLR_ACCENT   = RGBColor(0x00, 0xB4, 0xD8)   # Cyan accent
CLR_ACCENT2  = RGBColor(0x48, 0xCA, 0xA3)   # Teal green
CLR_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CLR_LIGHT    = RGBColor(0xC8, 0xE0, 0xF4)   # Light blue text
CLR_ORANGE   = RGBColor(0xFF, 0x9F, 0x1C)   # Orange highlight
CLR_RED      = RGBColor(0xFF, 0x4D, 0x6D)   # Red (bad ROI)
CLR_GREEN    = RGBColor(0x06, 0xD6, 0xA0)   # Green (good ROI)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

def fill_bg(slide, color=CLR_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, transparency=0):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if transparency:
        sh.fill.fore_color.theme_color = None
    sh.line.fill.background()
    return sh

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=CLR_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def add_notebook_ref(slide, notebook_name, section):
    """Thêm thanh tham chiếu notebook ở cuối slide"""
    add_rect(slide, 0.35, 6.75, 12.63, 0.58, RGBColor(0x0A, 0x1A, 0x2E))
    add_rect(slide, 0.35, 6.75, 0.06, 0.58, CLR_ACCENT2)  # accent strip
    add_text(slide, f"📓  Xem tại Notebook: ", 0.52, 6.78, 2.5, 0.45,
             size=11, bold=True, color=CLR_ACCENT2)
    add_text(slide, f"{notebook_name}  ", 3.0, 6.78, 5.5, 0.45,
             size=11, bold=True, color=CLR_ORANGE)
    add_text(slide, f"| Mục: {section}", 7.6, 6.78, 5.0, 0.45,
             size=11, bold=False, color=CLR_LIGHT)

def add_key_numbers(slide, numbers):
    """Thêm panel Key Numbers vào phần trống cuối slide
    numbers: list of (value, label) tuples, tối đa 4"""
    n = len(numbers)
    box_w = 12.63 / n
    for i, (val, label) in enumerate(numbers):
        bx = 0.35 + i * box_w
        add_rect(slide, bx, 6.0, box_w - 0.05, 0.68, RGBColor(0x13, 0x26, 0x42))
        add_rect(slide, bx, 6.0, box_w - 0.05, 0.07, CLR_ACCENT)  # top border
        add_text(slide, val, bx + 0.1, 6.02, box_w - 0.25, 0.35,
                 size=20, bold=True, color=CLR_ORANGE, align=PP_ALIGN.CENTER)
        add_text(slide, label, bx + 0.1, 6.38, box_w - 0.25, 0.28,
                 size=10, color=CLR_LIGHT, align=PP_ALIGN.CENTER)

def add_bullet_block(slide, items, l, t, w, h,
                      title=None, title_color=CLR_ACCENT,
                      size=15, indent_size=13):
    """items: list of (level, text, color) tuples"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True

    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(17)
        run.font.bold  = True
        run.font.color.rgb = title_color
        p = tf.add_paragraph()
        p.add_run().text = ""   # spacer

    for (level, text, color) in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        run.text = ("  " * level) + ("• " if level == 0 else "- ") + text
        run.font.size  = Pt(size if level == 0 else indent_size)
        run.font.color.rgb = color
    return tb

# =========================================================
# SLIDE 0 - TITLE
# =========================================================
slide = prs.slides.add_slide(BLANK)
fill_bg(slide)

# Top accent bar
add_rect(slide, 0, 0, 13.33, 0.07, CLR_ACCENT)

# Big title
add_text(slide,
         "Báo Cáo Kỹ Thuật: Quy Trình Sinh Dữ Liệu Giả Lập",
         0.6, 1.5, 12, 1.2,
         size=36, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(slide,
         "Synthetic Data Design & Structural Causal Model (SCM)",
         0.6, 2.8, 12, 0.6,
         size=22, bold=False, color=CLR_ACCENT, align=PP_ALIGN.CENTER)

# Meta info
add_text(slide,
         "Dự án: GSM Promotion Experimentation  |  Tuần 2  |  Intern: Data Science",
         0.6, 3.6, 12, 0.5,
         size=14, color=CLR_LIGHT, align=PP_ALIGN.CENTER)

# Bottom accent
add_rect(slide, 0, 7.43, 13.33, 0.07, CLR_ACCENT2)

# =========================================================
# SLIDE 1 - EXECUTIVE SUMMARY
# =========================================================
slide = prs.slides.add_slide(BLANK)
fill_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.06, CLR_ACCENT)

add_text(slide, "01  Executive Summary", 0.5, 0.12, 12, 0.55,
         size=26, bold=True, color=CLR_WHITE)

# Left box: Problem
add_rect(slide, 0.4, 0.9, 5.8, 5.9, RGBColor(0x13, 0x26, 0x42))
add_text(slide, "❓ Vấn đề", 0.6, 0.95, 5.4, 0.5,
         size=16, bold=True, color=CLR_ORANGE)
add_bullet_block(slide, [
    (0, "Dữ liệu quan sát (Observational Data) chứa Biến Nhiễu", CLR_LIGHT),
    (1, "Giờ cao điểm → ít phát Voucher nhưng khách vẫn đi nhiều", CLR_LIGHT),
    (1, "→ Gây ra Spurious Correlation (Tương quan ảo)", CLR_ORANGE),
    (0, "Không thể đo True ITE từ dữ liệu thực", CLR_LIGHT),
    (1, "Một người không thể vừa nhận vừa không nhận Voucher cùng lúc", CLR_LIGHT),
    (1, "→ Fundamental Problem of Causal Inference", CLR_ORANGE),
], 0.6, 1.5, 5.4, 4.8, size=14)

# Right box: Solution
add_rect(slide, 6.9, 0.9, 5.9, 5.9, RGBColor(0x13, 0x26, 0x42))
add_text(slide, "✅ Giải pháp: Structural Causal Model (SCM)", 7.1, 0.95, 5.5, 0.5,
         size=16, bold=True, color=CLR_GREEN)
add_bullet_block(slide, [
    (0, "Mô phỏng 20,000 khách hàng với Raw Features thực tế", CLR_LIGHT),
    (0, "Lập trình cứng (Hardcode) True ITE vào dữ liệu", CLR_LIGHT),
    (1, "→ Môi trường kiểm thử chuẩn mực cho A/B Testing", CLR_GREEN),
    (0, "3 thành phần thiết kế cốt lõi:", CLR_LIGHT),
    (1, "① Mô phỏng đặc điểm hành vi (Raw Features)", CLR_LIGHT),
    (1, "② Tích hợp Biến Nhiễu (Confounders)", CLR_LIGHT),
    (1, "③ Thiết lập Hiệu ứng Dị thể (HTE) theo Luật Nhân quả", CLR_LIGHT),
], 7.1, 1.5, 5.5, 4.8, size=14)

add_key_numbers(slide, [
    ("20,000",  "Khách hàng mô phỏng"),
    ("9",       "Raw Features thiết kế"),
    ("2",       "Biến Kết quả (y_obs, y_rand)"),
    ("1",       "Ground Truth (true_ite)"),
])
add_notebook_ref(slide,
    "2_complex_data_generation.ipynb",
    "Cell 1-2: Import & Tổng quan kiến trúc SCM")
add_rect(slide, 0, 7.43, 13.33, 0.07, CLR_ACCENT2)

# =========================================================
# SLIDE 2 - RAW FEATURES / DGP
# =========================================================
slide = prs.slides.add_slide(BLANK)
fill_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.06, CLR_ACCENT)
add_text(slide, "02  Mô Phỏng Hành Vi Khách Hàng (DGP)", 0.5, 0.12, 12, 0.55,
         size=26, bold=True, color=CLR_WHITE)

# Subtitle
add_text(slide, "Behavioral Proxies — không cần khai báo nhân khẩu học thực tế",
         0.5, 0.72, 12, 0.4, size=14, color=CLR_LIGHT)

rows = [
    ("is_urban",        "Bernoulli (p=0.7)",              "70% nội thành / 30% ngoại ô"),
    ("fare_obs",        "Log-normal",                     "Proxy cho Thu nhập. Cuốc sân bay × 4"),
    ("is_airport_trip", "Bernoulli (p=0.05)",             "5% khách. Tự động nhân cước × 4"),
    ("is_rush_hour",    "Phái sinh từ preferred_hour",    "Đi làm giờ 7-9h & 17-19h"),
    ("preferred_hour",  "Weighted Prob (EDA T1)",         "Đỉnh chóp vào giờ cao điểm"),
    ("is_rain_rider",   "Bernoulli (p=0.20)",             "Thích đi trời mưa → nhu cầu cao bất thường"),
    ("recency_days",    "Poisson",                        "Số ngày kể từ lần cuối đặt xe (0-30)"),
    ("monthly_rides",   "Poisson (tuỳ income & age)",     "Tần suất đi xe tháng trước"),
]

col_x = [0.35, 3.8, 7.4]
col_w = [3.3, 3.4, 5.5]
header_color = CLR_ACCENT

# Header
for ci, (cx, cw, label) in enumerate(zip(col_x, col_w,
        ["Tên biến", "Phân phối", "Ý nghĩa"])):
    add_rect(slide, cx, 1.2, cw, 0.45, RGBColor(0x02, 0x5E, 0x8C))
    add_text(slide, label, cx+0.05, 1.22, cw-0.1, 0.4,
             size=13, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)

for ri, (var, dist, desc) in enumerate(rows):
    row_y = 1.68 + ri * 0.65
    bg = RGBColor(0x13, 0x26, 0x42) if ri % 2 == 0 else RGBColor(0x19, 0x2F, 0x4E)
    for ci, (cx, cw) in enumerate(zip(col_x, col_w)):
        add_rect(slide, cx, row_y, cw, 0.62, bg)
    texts = [var, dist, desc]
    colors = [CLR_ACCENT, CLR_ORANGE, CLR_LIGHT]
    for ci, (cx, cw, txt, clr) in enumerate(zip(col_x, col_w, texts, colors)):
        add_text(slide, txt, cx+0.08, row_y+0.06, cw-0.15, 0.5,
                 size=12, color=clr, bold=(ci==0))

add_notebook_ref(slide,
    "2_complex_data_generation.ipynb",
    "Cell 3: UserSimulator.__init__() — xem phần phân phối từng biến")
add_rect(slide, 0, 7.43, 13.33, 0.07, CLR_ACCENT2)

# =========================================================
# SLIDE 3 - CONFOUNDERS
# =========================================================
slide = prs.slides.add_slide(BLANK)
fill_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.06, CLR_ACCENT)
add_text(slide, "03  Tích Hợp Biến Nhiễu (Confounders)", 0.5, 0.12, 12, 0.55,
         size=26, bold=True, color=CLR_WHITE)
add_text(slide,
         "Biến Nhiễu = Yếu tố tác động ĐỒNG THỜI lên Treatment (phát Voucher) VÀ Outcome (số chuyến)",
         0.5, 0.72, 12, 0.4, size=14, color=CLR_LIGHT)

# Definition box
add_rect(slide, 0.35, 1.18, 12.63, 0.85, RGBColor(0x02, 0x5E, 0x8C))
add_text(slide,
         "Nếu KHÔNG có Biến Nhiễu: Dữ liệu quá 'sạch' → Mô hình Uplift học không ra gì.\n"
         "Nếu CÓ Biến Nhiễu: Dữ liệu phức tạp y hệt đời thực → Kiểm tra được độ chính xác của Model.",
         0.6, 1.22, 12.2, 0.76, size=13, color=CLR_WHITE)

# Confounder 1
add_rect(slide, 0.35, 2.15, 6.0, 4.6, RGBColor(0x13, 0x26, 0x42))
add_text(slide, "🌆  is_rush_hour (Nhu cầu theo Giờ)",
         0.5, 2.2, 5.6, 0.45, size=16, bold=True, color=CLR_ORANGE)
add_bullet_block(slide, [
    (0, "Treatment ↓: Giờ cao điểm → App CẮT Voucher (tài xế khan hiếm)", CLR_LIGHT),
    (0, "Outcome ↑: Khách vẫn PHẢI đi làm dù không có mã giảm", CLR_LIGHT),
    (0, "Hậu quả nếu không kiểm soát:", CLR_ORANGE),
    (1, "Model nghĩ: 'Không Voucher → nhiều chuyến hơn' (NGƯỢC hoàn toàn!)", CLR_RED),
    (0, "Cách xử lý trong thí nghiệm:", CLR_ACCENT),
    (1, "Randomize (A/B Test) loại bỏ ảnh hưởng này", CLR_GREEN),
], 0.5, 2.7, 5.8, 4.0, size=13)

# Confounder 2
add_rect(slide, 7.0, 2.15, 6.0, 4.6, RGBColor(0x13, 0x26, 0x42))
add_text(slide, "🌧️  is_rain_rider (Cú sốc Thời tiết)",
         7.15, 2.2, 5.7, 0.45, size=16, bold=True, color=CLR_ORANGE)
add_bullet_block(slide, [
    (0, "Treatment ↓: Trời mưa → nền tảng TẮT Voucher (surge pricing)", CLR_LIGHT),
    (0, "Outcome ↑: Nhu cầu đặt xe tăng đột biến vì mưa", CLR_LIGHT),
    (0, "Hậu quả nếu không kiểm soát:", CLR_ORANGE),
    (1, "Thấy 'Ít Voucher → đi xe nhiều' → Rút kết luận SAI", CLR_RED),
    (0, "Cách xử lý trong thí nghiệm:", CLR_ACCENT),
    (1, "Đưa is_rain_rider vào làm Covariate trong mô hình", CLR_GREEN),
], 7.15, 2.7, 5.8, 4.0, size=13)

add_key_numbers(slide, [
    ("20%",   "Tỷ lệ Rain Riders trong dataset"),
    ("5%",    "Tỷ lệ Airport Business"),
    ("~30%",  "Tỷ lệ ngoại ô (Persuadables chính)"),
    ("~45%",  "Rush Hour users (Sure-things chính)"),
])
add_notebook_ref(slide,
    "2_complex_data_generation.ipynb",
    "Cell 4: _inject_confounders() — biến nhiễu tác động lên treatment_obs")
add_rect(slide, 0, 7.43, 13.33, 0.07, CLR_ACCENT2)

# =========================================================
# SLIDE 4 - HTE / CAUSAL RULES
# =========================================================
slide = prs.slides.add_slide(BLANK)
fill_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.06, CLR_ACCENT)
add_text(slide, "04  Thiết Lập Hiệu Ứng Dị Thể (HTE)", 0.5, 0.12, 12, 0.55,
         size=26, bold=True, color=CLR_WHITE)

# Formula
add_rect(slide, 0.35, 0.75, 12.63, 0.85, RGBColor(0x02, 0x5E, 0x8C))
add_text(slide,
         "ITE = 2.0 (Base)  −  1.5 × is_airport  −  1.0 × is_rush_hour  +  1.0 × (1 − is_urban)  +  0.5 × is_rain_rider",
         0.55, 0.78, 12.2, 0.76, size=14, bold=True, color=CLR_ORANGE, align=PP_ALIGN.CENTER)

# Left: Resistant
add_rect(slide, 0.35, 1.72, 6.0, 5.05, RGBColor(0x2A, 0x10, 0x14))
add_text(slide, "🔴  Nhóm Kháng Sale (Sure-things)", 0.5, 1.77, 5.7, 0.5,
         size=16, bold=True, color=CLR_RED)
add_bullet_block(slide, [
    (0, "Airport Business  →  ITE ≈ +0.5", CLR_LIGHT),
    (1, "Khách ra sân bay bắt buộc phải đi", CLR_LIGHT),
    (1, "Voucher là tiền vứt qua cửa sổ", CLR_RED),
    (0, "Urban Regulars (Rush Hour)  →  ITE ≈ +1.0", CLR_LIGHT),
    (1, "Đi làm hàng ngày, dù không có mã vẫn đi", CLR_LIGHT),
    (1, "Cannibalization Effect cực cao", CLR_RED),
    (0, "Rain Riders  →  ITE ≈ +1.5", CLR_LIGHT),
    (1, "Nhu cầu phát sinh từ thời tiết, không phải giá", CLR_LIGHT),
    (1, "Tặng Voucher không kích thêm được chuyến", CLR_RED),
], 0.5, 2.3, 5.7, 4.3, size=13)

# Right: Persuadable
add_rect(slide, 6.98, 1.72, 6.0, 5.05, RGBColor(0x08, 0x2A, 0x1E))
add_text(slide, "🟢  Nhóm Nhạy Cảm Giá (Persuadables)", 7.12, 1.77, 5.7, 0.5,
         size=16, bold=True, color=CLR_GREEN)
add_bullet_block(slide, [
    (0, "Suburban Occasionals  →  ITE ≈ +3.0", CLR_LIGHT),
    (1, "Khách ngoại ô, ít đi xe vì không có mã", CLR_LIGHT),
    (1, "Voucher thực sự kéo họ ra đường", CLR_GREEN),
    (0, "Urban Leisure  →  ITE ≈ +2.5", CLR_LIGHT),
    (1, "Khách đi chơi cuối tuần, nhạy cảm với giá", CLR_LIGHT),
    (1, "ROI dương trong kết quả A/B Test", CLR_GREEN),
    (0, "→ Voucher đúng là kích thích hành vi của nhóm này", CLR_ACCENT),
], 7.12, 2.3, 5.7, 4.3, size=13)

add_key_numbers(slide, [
    ("ITE = +0.5",  "Airport Business (kháng cao nhất)"),
    ("ITE = +1.0",  "Urban Regulars (đi làm)"),
    ("ITE = +2.5",  "Urban Leisure (nhạy cảm giá)"),
    ("ITE = +3.0",  "Suburban Occasionals (nhạy nhất)"),
])
add_notebook_ref(slide,
    "2_complex_data_generation.ipynb",
    "Cell 5: _compute_ite() — phương trình Causal Rules + true_ite")
add_rect(slide, 0, 7.43, 13.33, 0.07, CLR_ACCENT2)

# =========================================================
# SLIDE 5 - OUTCOME GENERATION
# =========================================================
slide = prs.slides.add_slide(BLANK)
fill_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.06, CLR_ACCENT)
add_text(slide, "05  Sinh Biến Kết Quả (Outcome Generation)", 0.5, 0.12, 12, 0.55,
         size=26, bold=True, color=CLR_WHITE)

# y_obs
add_rect(slide, 0.35, 0.8, 6.0, 3.1, RGBColor(0x13, 0x26, 0x42))
add_text(slide, "① y_obs  — Môi trường Tự nhiên", 0.55, 0.85, 5.6, 0.5,
         size=16, bold=True, color=CLR_ORANGE)
add_bullet_block(slide, [
    (0, "Số chuyến đi dựa HOÀN TOÀN vào đặc điểm khách hàng", CLR_LIGHT),
    (0, "Chưa có Voucher can thiệp", CLR_LIGHT),
    (0, "Bị nhiễu bởi treatment_obs (phát Voucher không ngẫu nhiên)", CLR_LIGHT),
    (0, "Dùng để demo Simpson's Paradox / Selection Bias", CLR_ORANGE),
], 0.55, 1.4, 5.7, 2.4, size=13)

# y_rand
add_rect(slide, 7.0, 0.8, 6.0, 3.1, RGBColor(0x13, 0x26, 0x42))
add_text(slide, "② y_rand  — Môi trường A/B Test (RCT)", 7.18, 0.85, 5.6, 0.5,
         size=16, bold=True, color=CLR_GREEN)
add_bullet_block(slide, [
    (0, "Tung đồng xu: 50% treatment_rand = 1", CLR_LIGHT),
    (0, "y_rand = y_obs + ITE × treatment  +  Poisson Noise", CLR_LIGHT),
    (0, "Poisson Distribution: đảm bảo số chuyến là số nguyên dương", CLR_LIGHT),
    (0, "→ Tiêu chuẩn vàng để đánh giá Uplift Model", CLR_GREEN),
], 7.18, 1.4, 5.7, 2.4, size=13)

# Ground Truth box
add_rect(slide, 0.35, 4.05, 12.63, 1.55, RGBColor(0x02, 0x5E, 0x8C))
add_text(slide, "⭐  true_ite  — Ground Truth (Điểm cốt lõi của Synthetic Data)", 0.55, 4.1, 12, 0.5,
         size=16, bold=True, color=CLR_ORANGE)
add_bullet_block(slide, [
    (0, "Là ITE thật sự của từng cá nhân — không bao giờ quan sát được trong dữ liệu thực", CLR_LIGHT),
    (0, "Lưu trữ trong dataset để kiểm chứng kết quả A/B Test: ROI thực tế vs. ITE thiết kế có khớp?", CLR_GREEN),
    (0, "→ LỢI THẾ độc nhất của Synthetic Data: có thể validate toàn bộ pipeline thống kê!", CLR_ORANGE),
], 0.55, 4.65, 12.1, 0.88, size=13)

add_key_numbers(slide, [
    ("y_obs",       "Môi trường quan sát (bị nhiễu)"),
    ("y_rand",      "RCT: phân bổ Voucher ngẫu nhiên"),
    ("true_ite",    "Ground Truth ITE từng cá nhân"),
    ("Poisson",     "Phân phối nhiễu ngẫu nhiên"),
])
add_notebook_ref(slide,
    "2_complex_data_generation.ipynb",
    "Cell 6-7: generate_outcomes() + visualize_distributions()")
add_rect(slide, 0, 7.43, 13.33, 0.07, CLR_ACCENT2)

# =========================================================
# SLIDE 6 - PIPELINE & NEXT STEPS
# =========================================================
slide = prs.slides.add_slide(BLANK)
fill_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.06, CLR_ACCENT)
add_text(slide, "06  Tính Liền Mạch Của Dự Án & Hướng Tiếp Theo", 0.5, 0.12, 12, 0.55,
         size=26, bold=True, color=CLR_WHITE)

stages = [
    ("Tuần 2", "Synthetic Data", "✔ Thiết kế SCM, gài luật HTE vào Raw Features, tạo Ground Truth", CLR_ACCENT),
    ("Tuần 3", "K-Means (5 cụm)", "✔ K-Means tự tìm ra 5 nhóm phản ánh đúng quy luật HTE", CLR_ACCENT2),
    ("Tuần 4", "A/B Test", "✔ ROI 5 nhóm khớp với luật nhân quả → Validate DGP thành công", CLR_ORANGE),
    ("Tuần 5", "AA Test", "✔ Kiểm tra phân bổ ngẫu nhiên hợp lệ (AA Sanity Check)", CLR_LIGHT),
    ("Tuần 6", "Báo cáo tổng kết", "✔ Tổng hợp toàn bộ pipeline, chuẩn bị cho hướng nghiên cứu tiếp theo", CLR_GREEN),
]
arrow_x = [0.28, 2.85, 5.42, 7.99, 10.56]
box_x   = [0.28, 2.85, 5.42, 7.99, 10.56]

for i, (week, name, desc, color) in enumerate(stages):
    bx = box_x[i]
    # box
    add_rect(slide, bx, 1.15, 2.4, 1.2, RGBColor(0x13, 0x26, 0x42))
    add_text(slide, week, bx+0.08, 1.18, 2.2, 0.38, size=12, bold=True, color=color)
    add_text(slide, name, bx+0.08, 1.56, 2.2, 0.72, size=14, bold=True, color=CLR_WHITE)
    # desc
    add_rect(slide, bx, 2.45, 2.4, 3.8, RGBColor(0x0A, 0x1A, 0x2E))
    add_text(slide, desc, bx+0.1, 2.52, 2.2, 3.6, size=12, color=CLR_LIGHT, wrap=True)
    # arrow (not after last)
    if i < 4:
        add_text(slide, "→", bx+2.42, 1.65, 0.4, 0.45, size=22, bold=True,
                 color=CLR_ACCENT, align=PP_ALIGN.CENTER)

# Conclusion
add_rect(slide, 0.35, 6.45, 12.63, 0.88, RGBColor(0x02, 0x5E, 0x8C))
add_text(slide,
         "Kết luận: Thiết kế DGP của Tuần 2 không phải sinh dữ liệu ngẫu nhiên — đây là một Phòng Thí Nghiệm "
         "Nhân Quả (Causal Laboratory) chuẩn mực, tạo nền tảng vững chắc cho A/B Testing và phân tích Causal Inference.",
         0.55, 6.5, 12.2, 0.78, size=13, color=CLR_WHITE)

add_notebook_ref(slide,
    "1_confounding_demo.ipynb  +  2_complex_data_generation.ipynb",
    "Chay tuan tu 2 notebooks de thay toan bo pipeline")
add_rect(slide, 0, 7.43, 13.33, 0.07, CLR_ACCENT2)

# ======================================
out_path = r"docs\Week2_Synthetic_Data_Deck_v2.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
