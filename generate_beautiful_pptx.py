import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Bảng màu phong cách Tech/AI chuyên nghiệp (Cảm hứng Xanh SM)
BG_COLOR = RGBColor(10, 25, 47) # Nền Xanh Đen (Dark Blue/Navy)
TITLE_COLOR = RGBColor(0, 242, 254) # Tiêu đề Xanh Cyan rực rỡ
TEXT_COLOR = RGBColor(230, 241, 255) # Chữ nội dung màu Trắng ngà
SUBTITLE_COLOR = RGBColor(100, 255, 218) # Phụ đề Xanh ngọc

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def format_title(shape, is_subtitle=False):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = SUBTITLE_COLOR if is_subtitle else TITLE_COLOR
            run.font.bold = True
            run.font.name = 'Segoe UI'

def format_body(shape):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = TEXT_COLOR
            run.font.name = 'Segoe UI'

def apply_theme(slide, is_title_slide=False):
    set_slide_background(slide)
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type == 1: # Center Title on Title Slide or Normal Title
                format_title(shape)
            elif shape.placeholder_format.type == 3: # Subtitle on Title Slide
                format_title(shape, is_subtitle=True)
            else:
                format_body(shape)

def create_week7_ppt():
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BÁO CÁO TUẦN 7: UPLIFT MODELING"
    slide.placeholders[1].text = "Tối ưu hóa Lợi nhuận bằng Causal AI"
    apply_theme(slide, is_title_slide=True)
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Tại sao cần Uplift Modeling?"
    tf = slide.placeholders[1].text_frame
    tf.text = "A/B Test truyền thống chỉ đo lường được kết quả trung bình (ATE) của đám đông."
    p = tf.add_paragraph()
    p.text = "Marketing cần cá nhân hóa: Đi tìm những người 'Bị thuyết phục' thay vì 'Đằng nào cũng mua'."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Uplift Modeling dự đoán chính xác CATE (Số chuyến đi tăng thêm sinh ra từ Voucher) cho TỪNG cá nhân."
    apply_theme(slide)
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mô hình T-Learner (XGBoost)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Thuật toán học máy kép (Two-Model Approach):"
    p = tf.add_paragraph()
    p.text = "Mô hình 1: Học từ nhóm Control (Để biết hành vi tự nhiên)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Mô hình 2: Học từ nhóm Treatment (Để biết hành vi khi có mã)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Điểm CATE = (Kết quả Mô hình 2) - (Kết quả Mô hình 1)."
    p = tf.add_paragraph()
    p.text = "Áp dụng Early Stopping và Validation Set để chống Overfitting."
    apply_theme(slide)
    
    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Bài toán Lợi nhuận & Điểm Hòa Vốn"
    tf = slide.placeholders[1].text_frame
    tf.text = "Bài toán Kinh tế cho nhóm Urban Cash:"
    p = tf.add_paragraph()
    p.text = "Lãi gộp = 20,000 VND | Chi phí Voucher = 15,000 VND"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Ngưỡng CATE hòa vốn = 15k / 20k = 0.75 chuyến"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Thực tế: AI phát hiện KHÔNG AI trong tập Urban Cash vượt qua ngưỡng này."
    p = tf.add_paragraph()
    p.text = "Kết luận: DỪNG CHIẾN DỊCH, tránh lỗ hàng triệu đồng so với phát đại trà."
    apply_theme(slide)
    
    prs.save('docs/Week7_Uplift_Modeling_Summary.pptx')

def create_week8_ppt():
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BÁO CÁO TUẦN 8: STRESS TESTING"
    slide.placeholders[1].text = "Kiểm định Tính vững (Robustness Check)"
    apply_theme(slide, is_title_slide=True)
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mục tiêu của Stress Test"
    tf = slide.placeholders[1].text_frame
    tf.text = "Bảo vệ hệ thống phân tích trước những biến động khó lường của thực tế."
    p = tf.add_paragraph()
    p.text = "Trả lời câu hỏi C-Level: Mô hình có bị gãy nếu dữ liệu bị nhiễu hoặc có quy mô quá lớn?"
    apply_theme(slide)
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4 Bài Test Khắc Nghiệt"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Scale-up: Tăng mẫu lên 100,000 người -> ATE ổn định, P-value cực nhỏ."
    p = tf.add_paragraph()
    p.text = "2. A/A Test: Giả lập Voucher hỏng (Effect=0) -> Báo không có tác dụng (Không False Positive)."
    p = tf.add_paragraph()
    p.text = "3. Lệch tỷ lệ chia mẫu (90/10): Ngân sách hẹp -> ATE vẫn bảo toàn."
    p = tf.add_paragraph()
    p.text = "4. Bơm nhiễu (Gaussian Noise): Bão lụt, kẹt xe -> Randomization triệt tiêu nhiễu hoàn toàn."
    apply_theme(slide)
    
    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Khẳng định Thành quả Dự án"
    tf = slide.placeholders[1].text_frame
    tf.text = "Khung phân tích Causal Inference của dự án là hoàn toàn Vững (Robust)."
    p = tf.add_paragraph()
    p.text = "Sẵn sàng Deploy lên hệ thống Production của Xanh SM."
    p = tf.add_paragraph()
    p.text = "Mô hình Uplift sẵn sàng đóng vai trò 'Người gác cổng' tối ưu ROI tự động."
    apply_theme(slide)
    
    prs.save('docs/Week8_Stress_Test_Summary.pptx')

if __name__ == '__main__':
    create_week7_ppt()
    create_week8_ppt()
