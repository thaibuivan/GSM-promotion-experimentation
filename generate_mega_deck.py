import sys
import os

sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from generate_awesome_pptx import (
    create_title_slide, create_base_layout, add_cards,
    YELLOW, CYAN, GREEN, PINK, WHITE
)

def create_mega_ppt():
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    TOTAL = 10
    current_slide = 1
    
    create_title_slide(prs, "BÁO CÁO TIẾN ĐỘ: TỪ DATA ĐẾN A/B TESTING", "Tổng hợp: Mô phỏng Dữ liệu, K-Means & Thử nghiệm")
    
    # --- PART 1: SYNTHETIC DATA ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s1, f"{current_slide}. Xây dựng Sandbox Dữ liệu (DGP)", "Kế thừa từ EDA Tuần 1", current_slide, TOTAL)
    add_cards(s1, [
        {'title': 'Vấn đề Dữ liệu thực', 'content': 'Dữ liệu lịch sử chứa nhiều Bias và tương quan ảo (VD: Giờ cao điểm ít voucher nhưng nhiều chuyến). Không thể dùng để test Model.', 'accent': PINK},
        {'title': 'Giải pháp Mô phỏng', 'content': 'Tạo ra 20,000 khách hàng ảo với thông số kế thừa từ EDA thật (Avg Fare $17.60). Đây là phòng thí nghiệm hoàn hảo có sẵn Ground Truth.', 'accent': GREEN}
    ])
    current_slide += 1

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s2, f"{current_slide}. Tích hợp Biến nhiễu & HTE", "Tạo ra độ khó cho Mô hình", current_slide, TOTAL)
    add_cards(s2, [
        {'title': 'Biến Nhiễu (Confounders)', 'content': 'Đưa yếu tố Giờ cao điểm và Thời tiết vào để tạo ra các tương quan phức tạp như ngoài đời thực.', 'accent': YELLOW},
        {'title': 'Hiệu ứng Dị thể (HTE)', 'content': 'Chia khách thành nhóm "Sure-Things" (đằng nào cũng đi) và "Persuadables" (nhạy cảm giá). Mỗi nhóm phản ứng với Voucher một cách khác nhau.', 'accent': CYAN}
    ])
    current_slide += 1

    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s3, f"{current_slide}. Ground Truth (True ITE)", "Vũ khí bí mật của Synthetic Data", current_slide, TOTAL)
    add_cards(s3, [
        {'title': 'True ITE', 'content': 'Vì chúng ta tự tạo ra dữ liệu, chúng ta biết chính xác Tác động nhân quả thực sự của từng khách hàng (Thứ không bao giờ thấy ở đời thực).', 'accent': GREEN},
        {'title': 'Giá trị cốt lõi', 'content': 'True ITE sẽ dùng làm đáp án chuẩn (Benchmark) để kiểm định độ chính xác của các thuật toán Uplift ở những tuần sau.', 'accent': CYAN}
    ])
    current_slide += 1

    # --- PART 2: WEEK 3 CONTENT (SEGMENTATION) ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s4, f"{current_slide}. Phân cụm Khách hàng", "Lựa chọn Tập mục tiêu bằng K-Means", current_slide, TOTAL)
    add_cards(s4, [
        {'title': 'Tiêu chí Chọn vào (Inclusion)', 'content': 'Thay vì phát mass, chỉ nhắm vào nhóm Persuadables (Nhạy cảm giá): Urban Credit Card và Suburban Occasionals.', 'accent': CYAN},
        {'title': 'Tiêu chí Loại trừ (Exclusion)', 'content': 'Khách đi sân bay (Kháng sale) và Khách đi làm giờ cao điểm (Bắt buộc phải đi).', 'accent': PINK}
    ])
    current_slide += 1
    
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s5, f"{current_slide}. Thiết kế Can thiệp", "Treatment vs Control", current_slide, TOTAL)
    add_cards(s5, [
        {'title': 'Nhóm Can thiệp (Treatment)', 'content': 'Gửi Push Notification, In-app Banner và tự động áp mã giảm 25% khi thanh toán.', 'accent': GREEN},
        {'title': 'Nhóm Đối chứng (Control)', 'content': 'Không nhận bất cứ thông báo nào. Không chạy chiến dịch đè lên nhau.', 'accent': YELLOW}
    ])
    current_slide += 1
    
    # --- PART 3: WEEK 4 CONTENT (A/B TESTING) ---
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s6, f"{current_slide}. Hệ thống Chỉ số", "OEC & Guardrails", current_slide, TOTAL)
    add_cards(s6, [
        {'title': 'Chỉ số Cốt lõi (OEC)', 'content': 'Incremental Rides (Số chuyến đi tăng thêm nhờ Voucher).', 'accent': CYAN},
        {'title': 'Chỉ số Bảo vệ (Guardrails)', 'content': 'Profit Margin (Lợi nhuận không âm). Đảm bảo không bù lỗ vô tội vạ trên từng cuốc xe.', 'accent': YELLOW}
    ])
    current_slide += 1
    
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s7, f"{current_slide}. Nền tảng Thống kê", "Kiểm định Giả thuyết (Frequentist Testing)", current_slide, TOTAL)
    add_cards(s7, [
        {'title': 'Giả thuyết Không (H0)', 'content': 'Giả định rằng Voucher không tạo ra sự khác biệt nào giữa 2 nhóm.', 'accent': YELLOW},
        {'title': 'P-value', 'content': 'Đại diện cho xác suất thu được chênh lệch do ngẫu nhiên.\nNgưỡng Alpha = 0.05 là tiêu chuẩn để bác bỏ H0.', 'accent': CYAN}
    ])
    current_slide += 1
    
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s8, f"{current_slide}. Hai Tiêu chí Đánh giá Độc lập", "Sanity vs OEC", current_slide, TOTAL)
    add_cards(s8, [
        {'title': 'Sanity Checks', 'content': 'Kiểm tra độ cân bằng hệ thống phân bổ.\nKỳ vọng P-value > 0.05 để chứng minh hệ thống không bị lỗi.', 'accent': PINK},
        {'title': 'OEC Evaluation', 'content': 'Đo lường tác động lên chỉ số cốt lõi.\nKỳ vọng P-value < 0.05 để chứng minh Voucher có hiệu quả thật sự.', 'accent': GREEN}
    ])
    current_slide += 1
    
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s9, f"{current_slide}. Đánh giá Tác động Ban đầu", "Chênh lệch giữa Control và Treatment", current_slide, TOTAL)
    add_cards(s9, [
        {'title': 'Quan sát trực quan', 'content': 'Nhóm nhận khuyến mãi (Treatment) có sự dịch chuyển tích cực về số chuyến đi.', 'accent': GREEN},
        {'title': 'Khẳng định', 'content': 'Xác nhận Voucher thực sự tạo ra Tác động Nhân quả (Causal Effect) ở mức độ tổng thể.', 'accent': CYAN}
    ])
    current_slide += 1

    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    create_base_layout(s10, f"{current_slide}. Kế hoạch Tiếp theo", "Từ A/B Testing đến Uplift Modeling", current_slide, TOTAL)
    add_cards(s10, [
        {'title': 'Giới hạn của A/B Test', 'content': 'Chỉ trả lời được câu hỏi "Chiến dịch có hiệu quả chung không?", không tối ưu được cho từng cá nhân cụ thể.', 'accent': YELLOW},
        {'title': 'Mục tiêu Tuần tới', 'content': 'Huấn luyện Machine Learning (Uplift Model) để trả lời: Đưa Voucher cho khách hàng A hay khách hàng B thì sẽ thu được lợi nhuận cao nhất?', 'accent': CYAN}
    ])

    prs.save('docs/Progress_Update_Mega_Deck.pptx')
    print("Done generating docs/Progress_Update_Mega_Deck.pptx")

if __name__ == '__main__':
    create_mega_ppt()
