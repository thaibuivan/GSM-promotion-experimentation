import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_week7_ppt():
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Báo cáo Tuần 7: Uplift Modeling"
    subtitle.text = "Tối ưu hóa Lợi nhuận bằng Causal AI"
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Tại sao cần Uplift Modeling?"
    tf = slide.placeholders[1].text_frame
    tf.text = "A/B Test truyền thống chỉ đo lường được kết quả trung bình (ATE) của đám đông."
    p = tf.add_paragraph()
    p.text = "Tuy nhiên, Marketing cần sự cá nhân hóa: Đi tìm những người 'Bị thuyết phục' thay vì những người 'Đằng nào cũng mua' (Sure Things)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Uplift Modeling dự đoán chính xác CATE (Số chuyến đi tăng thêm sinh ra từ cái Voucher) cho TỪNG người."
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mô hình T-Learner (XGBoost)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Thuật toán học máy kép (Two-Model Approach):"
    p = tf.add_paragraph()
    p.text = "Mô hình 1: Học từ nhóm Control (Để biết hành vi tự nhiên)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Mô hình 2: Học từ nhóm Treatment (Để biết hành vi khi có mã)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Điểm CATE = (Kết quả Mô hình 2) - (Kết quả Mô hình 1)."
    p = tf.add_paragraph()
    p.text = "Đã áp dụng Early Stopping và Validation Set (60-20-20) để chống Overfitting."
    
    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Bài toán Lợi nhuận & Điểm Hòa Vốn"
    tf = slide.placeholders[1].text_frame
    tf.text = "Bài toán Kinh tế cho nhóm Urban Cash:"
    p = tf.add_paragraph()
    p.text = "Lãi gộp = 20,000 VND | Chi phí Voucher = 15,000 VND"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Ngưỡng CATE hòa vốn = 15k / 20k = 0.75 chuyến"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Thực tế: Quét qua tập Urban Cash, mô hình phát hiện KHÔNG AI vượt qua ngưỡng này."
    p = tf.add_paragraph()
    p.text = "Kết luận: AI quyết định DỪNG CHIẾN DỊCH, giúp công ty tránh lỗ hàng triệu đồng so với việc phát đại trà."
    
    prs.save('docs/Week7_Uplift_Modeling_Summary.pptx')

def create_week8_ppt():
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Báo cáo Tuần 8: Stress Testing"
    subtitle.text = "Kiểm định Tính vững (Robustness Check)"
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mục tiêu của Stress Test"
    tf = slide.placeholders[1].text_frame
    tf.text = "Bảo vệ hệ thống phân tích trước những biến động khó lường của thực tế."
    p = tf.add_paragraph()
    p.text = "Trả lời câu hỏi: Mô hình có bị gãy nếu dữ liệu thực tế bị nhiễu hoặc có quy mô quá lớn?"
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4 Bài Test Khắc Nghiệt"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Scale-up: Tăng mẫu lên 100,000 người -> ATE ổn định, P-value cực nhỏ."
    p = tf.add_paragraph()
    p.text = "2. A/A Test: Giả lập Voucher hỏng (Effect=0) -> Hệ thống báo không có tác dụng (Không False Positive)."
    p = tf.add_paragraph()
    p.text = "3. Lệch tỷ lệ chia mẫu (90/10): Ngân sách hẹp -> ATE vẫn bảo toàn."
    p = tf.add_paragraph()
    p.text = "4. Bơm nhiễu (Gaussian Noise): Bão lụt, kẹt xe -> Randomization triệt tiêu nhiễu tuyệt đối."
    
    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Khẳng định Thành quả Dự án"
    tf = slide.placeholders[1].text_frame
    tf.text = "Khung phân tích Causal Inference của chúng ta hoàn toàn Vững (Robust)."
    p = tf.add_paragraph()
    p.text = "Sẵn sàng Deploy lên hệ thống Production của Xanh SM."
    p = tf.add_paragraph()
    p.text = "Mô hình Uplift sẵn sàng đóng vai trò 'Người gác cổng' tiết kiệm dòng tiền."
    
    prs.save('docs/Week8_Stress_Test_Summary.pptx')

if __name__ == '__main__':
    create_week7_ppt()
    create_week8_ppt()
    print("Đã tạo thành công 2 file PowerPoint!")
